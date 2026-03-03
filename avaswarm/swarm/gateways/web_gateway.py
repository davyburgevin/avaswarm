"""
Web gateway — FastAPI multi-page UI with:
  - /            Home (cards)
  - /chat        Chat page (WebSocket streaming)
  - /config      LLM configuration
  - /memory      Memory browser
  - /scheduler   Job manager

REST API under /api/*
Static files from swarm/web/static/
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Any

import uvicorn
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, Request, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
import json

from swarm.config import settings
from swarm.gateways.base import BaseGateway
import httpx

if TYPE_CHECKING:
    from swarm.core.agent import Agent

logger = logging.getLogger(__name__)

_WEB_DIR       = Path(__file__).resolve().parent.parent / "web"
_TEMPLATES_DIR = _WEB_DIR / "templates"
_STATIC_DIR    = _WEB_DIR / "static"

# ---------------------------------------------------------------------------
# Helper — persist _runtime_config to .env so settings survive restarts
# ---------------------------------------------------------------------------
_ENV_FILE = Path(__file__).resolve().parent.parent.parent / ".env"

# Map ConfigPayload field → .env key name
_ENV_KEYS = {
    "provider":             "DEFAULT_PROVIDER",
    "model":                "DEFAULT_MODEL",
    "openai_api_key":       "OPENAI_API_KEY",
    "openai_base_url":      "OPENAI_BASE_URL",
    "anthropic_api_key":    "ANTHROPIC_API_KEY",
    "openrouter_api_key":   "OPENROUTER_API_KEY",
    "openrouter_base_url":  "OPENROUTER_BASE_URL",
    "github_copilot_token": "GITHUB_COPILOT_TOKEN",
    "github_oauth_client_id": "GITHUB_OAUTH_CLIENT_ID",
    "vllm_base_url":        "VLLM_BASE_URL",
    "context_window":       "CONTEXT_WINDOW",
}

def _persist_env(data: dict) -> None:
    """Write (or update) key=value lines in .env without touching unknown lines."""
    # Read existing lines if file exists
    existing: dict[str, str] = {}
    lines: list[str] = []
    if _ENV_FILE.exists():
        for raw in _ENV_FILE.read_text(encoding="utf-8").splitlines():
            line = raw.strip()
            if line and not line.startswith("#") and "=" in line:
                k, _, v = line.partition("=")
                existing[k.strip()] = v.strip()
            lines.append(raw)

    # Merge new values
    for payload_key, env_key in _ENV_KEYS.items():
        value = data.get(payload_key)
        if value:  # only persist non-empty values
            existing[env_key] = value

    # Rebuild file: update in-place lines, append new ones
    updated_keys: set[str] = set()
    new_lines: list[str] = []
    for raw in lines:
        line = raw.strip()
        if line and not line.startswith("#") and "=" in line:
            k, _, _ = line.partition("=")
            k = k.strip()
            if k in existing:
                new_lines.append(f"{k}={existing[k]}")
                updated_keys.add(k)
                continue
        new_lines.append(raw)

    for env_key, value in existing.items():
        if env_key not in updated_keys:
            new_lines.append(f"{env_key}={value}")

    _ENV_FILE.write_text("\n".join(new_lines) + "\n", encoding="utf-8")
    logger.info("Config persisted to %s", _ENV_FILE)


# ---------------------------------------------------------------------------
# Multiple named LLM configs persistence
# ---------------------------------------------------------------------------
_CONFIGS_FILE = Path(__file__).resolve().parent.parent.parent / "llm_configs.json"

def _load_configs() -> dict:
    if not _CONFIGS_FILE.exists():
        return {"configs": [], "active": None}
    try:
        import json
        return json.loads(_CONFIGS_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {"configs": [], "active": None}

def _save_configs(data: dict) -> None:
    import json
    _CONFIGS_FILE.write_text(json.dumps(data, indent=2), encoding="utf-8")


# ---------------------------------------------------------------------------
# Helper — build a live provider instance from a ConfigPayload
# ---------------------------------------------------------------------------
def _provider_from_payload(payload: "ConfigPayload"):
    """Instantiate the right provider with the correct kwargs from the form."""
    from swarm.providers.registry import registry
    provider_cls = registry.get(payload.provider)
    kwargs: dict = {}
    if payload.provider == "openai":
        if payload.openai_api_key:
            kwargs["api_key"] = payload.openai_api_key
        if payload.openai_base_url:
            kwargs["base_url"] = payload.openai_base_url
    elif payload.provider == "anthropic":
        if payload.anthropic_api_key:
            kwargs["api_key"] = payload.anthropic_api_key
    elif payload.provider == "openrouter":
        if payload.openrouter_api_key:
            kwargs["api_key"] = payload.openrouter_api_key
        if payload.openrouter_base_url:
            kwargs["base_url"] = payload.openrouter_base_url
    elif payload.provider == "github_copilot":
        if payload.github_copilot_token:
            kwargs["oauth_token"] = payload.github_copilot_token
    elif payload.provider == "vllm":
        if payload.vllm_base_url:
            kwargs["base_url"] = payload.vllm_base_url
        kwargs["model"] = payload.model  # critical — sets _default_model
    return provider_cls(**kwargs)


# Runtime config (hot-swappable, survives for the process lifetime)
_runtime_config: dict = {
    "provider":             settings.default_provider,
    "model":                settings.default_model,
    "openai_api_key":       settings.openai_api_key,
    "openai_base_url":      settings.openai_base_url,
    "anthropic_api_key":    settings.anthropic_api_key,
    "openrouter_api_key":   settings.openrouter_api_key,
    "openrouter_base_url":  settings.openrouter_base_url,
    "github_copilot_token": settings.github_copilot_token,
    "github_oauth_client_id": settings.github_oauth_client_id,
    "vllm_base_url":        settings.vllm_base_url,
    "context_window":       settings.context_window,
}


# ---- Pydantic models ----

class ConfigPayload(BaseModel):
    provider: str
    model: str
    name: str = ""
    openai_api_key: str = ""
    openai_base_url: str = ""
    anthropic_api_key: str = ""
    openrouter_api_key: str = ""
    openrouter_base_url: str = ""
    github_copilot_token: str = ""
    github_oauth_client_id: str = ""
    vllm_base_url: str = ""
    context_window: int = 16000

class JobPayload(BaseModel):
    id: str
    cron: str
    task: str

class MemoryPayload(BaseModel):
    content: str


# ---- Gateway ----

class WebGateway(BaseGateway):
    name = "web"

    def __init__(self, agent: "Agent", host: str | None = None, port: int | None = None) -> None:
        super().__init__(agent)
        self._host = host or settings.web_host
        self._port = port or settings.web_port
        self._server: uvicorn.Server | None = None
        self.app = self._build_app()

    def _build_app(self) -> FastAPI:
        app = FastAPI(title="Swarm", version="0.1.0")
        app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

        templates = Jinja2Templates(directory=str(_TEMPLATES_DIR))
        app.mount("/static", StaticFiles(directory=str(_STATIC_DIR)), name="static")

        agent = self.agent

        # ==============================================================
        # PAGE ROUTES
        # ==============================================================

        @app.get("/", response_class=HTMLResponse)
        async def page_home(request: Request):
            # Try to extract first name from long-term memory (saved as 'UserFirstName: NAME')
            first_name = None
            try:
                lt = agent.memory.get_long_term_sync() or ""
                for line in lt.splitlines():
                    if line.startswith("UserFirstName:"):
                        first_name = line.split(":", 1)[1].strip()
                        break
            except Exception:
                first_name = None
            return templates.TemplateResponse("index.html", {"request": request, "first_name": first_name})

        @app.get("/chat", response_class=HTMLResponse)
        async def page_chat(request: Request):
            return templates.TemplateResponse("chat.html", {
                "request": request,
                "model": _runtime_config["model"],
            })

        @app.get("/config", response_class=HTMLResponse)
        async def page_config(request: Request):
            cfgs = _load_configs()
            return templates.TemplateResponse("config.html", {
                "request": request,
                "current": _runtime_config,
                "configs": cfgs.get("configs", []),
                "active": cfgs.get("active"),
            })

        @app.get("/memory", response_class=HTMLResponse)
        async def page_memory(request: Request):
            from datetime import date
            lt    = await agent.memory.read_long_term()
            daily = await agent.memory.read_daily()
            return templates.TemplateResponse("memory.html", {
                "request":   request,
                "long_term": lt,
                "daily":     daily,
                "today":     date.today().isoformat(),
            })

        @app.get("/scheduler", response_class=HTMLResponse)
        async def page_scheduler(request: Request):
            from swarm.scheduler import scheduler
            return templates.TemplateResponse("scheduler.html", {
                "request": request,
                "jobs":    scheduler.list_jobs(),
            })

        @app.get("/agents", response_class=HTMLResponse)
        async def page_agents(request: Request):
            # List existing sub-agents from the project 'agents' directory
            agents_dir = Path(__file__).resolve().parent.parent.parent / "agents"
            agents: list[str] = []
            try:
                if agents_dir.exists():
                    for p in sorted(agents_dir.iterdir()):
                        if p.is_dir():
                            agents.append(p.name)
            except Exception:
                agents = []
            return templates.TemplateResponse("agents.html", {"request": request, "agents": agents})

        @app.get("/agents/new", response_class=HTMLResponse)
        async def page_create_agent(request: Request):
            return templates.TemplateResponse("create_agent.html", {"request": request})

        @app.get("/agents/{name}/instructions", response_class=HTMLResponse)
        async def page_agent_instructions(request: Request, name: str):
            target = _agents_dir / name
            if not target.is_dir():
                from fastapi.responses import RedirectResponse
                return RedirectResponse("/agents")
            instr_path = target / "agent.md"
            content = instr_path.read_text("utf-8") if instr_path.exists() else _DEFAULT_AGENT_INSTRUCTIONS
            return templates.TemplateResponse("agent_instructions.html", {
                "request": request,
                "agent_name": name,
                "content": content,
            })

        # ==============================================================
        # API — CONFIG
        # ==============================================================

        @app.get("/api/config")
        async def api_get_config():
            cfgs = _load_configs()
            base = {k: ("***" if ("key" in k or "token" in k) and v else v)
                for k, v in _runtime_config.items()}
            base["configs"] = cfgs.get("configs", [])
            base["active"] = cfgs.get("active")
            return base

        @app.post("/api/config/test")
        async def api_test_config(payload: ConfigPayload):
            """Test a provider configuration by instantiating the provider and
            sending a tiny completion request. Returns logs and any error.
            """
            logs: list[str] = []
            def log(msg: str) -> None:
                logs.append(str(msg))

            try:
                log(f"Provider: {payload.provider} | Model: {payload.model}")

                # For Copilot, pre-check that the model exists
                if payload.provider == "github_copilot":
                    try:
                        # Use token from runtime config (updated by device login) or settings
                        _tok = (
                            payload.github_copilot_token
                            or _runtime_config.get("github_copilot_token")
                            or settings.github_copilot_token
                        )
                        if not _tok:
                            raise RuntimeError("No token")
                        # Try internal exchange; fall back to direct bearer
                        _bearer = _tok
                        try:
                            async with httpx.AsyncClient() as _c:
                                _ex = await _c.get(
                                    "https://api.github.com/copilot_internal/v2/token",
                                    headers={
                                        "Authorization": f"token {_tok}",
                                        "Accept": "application/json",
                                        "User-Agent": "GithubCopilot/1.246.0",
                                        "Editor-Version": "vscode/1.96.0",
                                        "Editor-Plugin-Version": "copilot/1.246.0",
                                    },
                                    timeout=10,
                                )
                                if _ex.is_success:
                                    _bearer = _ex.json().get("token", _tok)
                        except Exception:
                            pass
                        async with httpx.AsyncClient() as _c:
                            _mr = await _c.get(
                                "https://api.githubcopilot.com/models",
                                headers={
                                    "Authorization": f"Bearer {_bearer}",
                                    "Editor-Version": "vscode/1.96.0",
                                    "Copilot-Integration-Id": "vscode-chat",
                                    "User-Agent": "GithubCopilot/1.246.0",
                                },
                                timeout=10,
                            )
                        if _mr.is_success:
                            _available = [m["id"] for m in _mr.json().get("data", [])]
                            log(f"Available models ({len(_available)}): {', '.join(_available)}")
                            if payload.model and payload.model not in _available:
                                log(f"⚠ Model '{payload.model}' is NOT in the available list!")
                                return {
                                    "ok": False,
                                    "error": f"Model '{payload.model}' is not available on your Copilot plan.",
                                    "detail": f"Available models: {', '.join(_available)}",
                                    "logs": logs,
                                }
                    except Exception as _me:
                        log(f"Model pre-check skipped: {_me}")

                # vLLM may require a lightweight HTTP probe before the full test
                if payload.provider == "vllm":
                    url = payload.vllm_base_url or settings.vllm_base_url
                    try:
                        base = url.rstrip('/')
                        base_no_v1 = base.rsplit('/v1', 1)[0]
                        async with httpx.AsyncClient() as c:
                            ping = await c.get(base_no_v1 + '/api/tags', timeout=3)
                        log(f"vLLM /api/tags HTTP {ping.status_code}")
                    except Exception as ping_err:
                        log(f"vLLM reachability probe failed: {ping_err}")

                # Instantiate provider and send a tiny completion
                prov = _provider_from_payload(payload)
                log("Sending test completion request...")
                from swarm.providers.base import CompletionRequest, Message
                req = CompletionRequest(
                    messages=[Message(role="user", content="Reply with only: OK")],
                    model=payload.model,
                    max_tokens=5,
                )
                resp = await prov.complete(req)
                log(f"Success! model={resp.model} reply={resp.content[:40]!r}")
                return {"ok": True, "model": resp.model, "reply": resp.content[:80], "logs": logs}
            except Exception as exc:
                import traceback
                tb = traceback.format_exc()
                log(f"ERROR: {exc}")
                logger.error("[config/test] %s\n%s", exc, tb)
                return {"ok": False, "error": str(exc), "detail": tb, "logs": logs}

        # =============================================================
        # Copilot device-login helpers
        # =============================================================

        @app.post('/api/copilot/device/start')
        async def api_copilot_device_start(body: dict | None = None):
            """Start GitHub OAuth device flow and return user_code + verification_uri.

            The client_id may be provided in the POST body as {"client_id": "..."}.
            If not provided we fall back to runtime config and then the .env settings.
            """
            client_id = None
            if body and isinstance(body, dict):
                client_id = body.get("client_id")
            if not client_id:
                client_id = _runtime_config.get("github_oauth_client_id") or settings.github_oauth_client_id
            if not client_id:
                return {"ok": False, "error": "GITHUB_OAUTH_CLIENT_ID not configured (provide it in the form or .env)"}
            try:
                async with httpx.AsyncClient() as c:
                    r = await c.post(
                        'https://github.com/login/device/code',
                        data={"client_id": client_id, "scope": "read:user copilot"},
                        headers={"Accept": "application/json"},
                        timeout=10,
                    )
                r.raise_for_status()
                return {"ok": True, **r.json()}
            except Exception as exc:
                logger.error("device/start error: %s", exc)
                return {"ok": False, "error": str(exc)}

        @app.post('/api/copilot/device/exchange')
        async def api_copilot_device_exchange(body: dict):
            """Attempt one exchange of device_code for an access token. Client should poll repeatedly.

            Returns GitHub token on success. On pending/slow_down/authorization_pending, return that status.
            """
            device_code = body.get('device_code')
            # allow client to pass client_id directly; otherwise use runtime config or settings
            client_id = body.get('client_id') or _runtime_config.get('github_oauth_client_id') or settings.github_oauth_client_id
            if not client_id or not device_code:
                return {"ok": False, "error": "Missing client_id or device_code"}
            try:
                async with httpx.AsyncClient() as c:
                    r = await c.post(
                        'https://github.com/login/oauth/access_token',
                        data={
                            'client_id': client_id,
                            'device_code': device_code,
                            'grant_type': 'urn:ietf:params:oauth:grant-type:device_code',
                        },
                        headers={"Accept": "application/json"},
                        timeout=10,
                    )
                r.raise_for_status()
                data = r.json()
                # Possible keys: access_token, error, error_description
                if data.get('access_token'):
                    token = data['access_token']
                    # Persist token to .env using payload key so _persist_env maps it
                    _persist_env({'github_copilot_token': token})
                    # Also update runtime config so live agent uses it
                    _runtime_config['github_copilot_token'] = token
                    try:
                        # If agent exists in closure, set provider token if active
                        if agent and getattr(agent, 'provider', None) and getattr(agent.provider, 'name', '') == 'github_copilot':
                            agent.provider._oauth_token = token  # type: ignore[attr-defined]
                            agent.provider._session_token = None  # type: ignore[attr-defined]
                        # Clear module-level session cache so next call re-runs exchange
                        try:
                            from swarm.providers.github_copilot_provider import _SESSION_CACHE
                            _SESSION_CACHE.clear()
                        except Exception:
                            pass
                    except Exception:
                        pass
                    return {"ok": True, "access_token": token}
                # Return error info (e.g., authorization_pending, slow_down)
                return {"ok": False, "status": data.get('error'), "detail": data.get('error_description')}
            except Exception as exc:
                logger.error('device/exchange error: %s', exc)
                return {"ok": False, "error": str(exc)}

        @app.get('/api/copilot/models')
        async def api_copilot_models():
            """Fetch the live model list from api.githubcopilot.com using the stored token."""
            token = _runtime_config.get('github_copilot_token') or settings.github_copilot_token
            if not token:
                return {"ok": False, "error": "No GitHub Copilot token configured", "models": []}
            # Resolve bearer token (try internal exchange first, fall back to OAuth token directly)
            bearer = token
            try:
                async with httpx.AsyncClient() as c:
                    r = await c.get(
                        'https://api.github.com/copilot_internal/v2/token',
                        headers={
                            "Authorization": f"token {token}",
                            "Accept": "application/json",
                            "User-Agent": "GithubCopilot/1.246.0",
                            "Editor-Version": "vscode/1.96.0",
                            "Editor-Plugin-Version": "copilot/1.246.0",
                        },
                        timeout=10,
                    )
                    if r.is_success:
                        bearer = r.json().get('token', token)
            except Exception:
                pass
            try:
                async with httpx.AsyncClient() as c:
                    r = await c.get(
                        'https://api.githubcopilot.com/models',
                        headers={
                            "Authorization": f"Bearer {bearer}",
                            "Content-Type": "application/json",
                            "Editor-Version": "vscode/1.96.0",
                            "Copilot-Integration-Id": "vscode-chat",
                            "User-Agent": "GithubCopilot/1.246.0",
                        },
                        timeout=15,
                    )
                r.raise_for_status()
                data = r.json()
                models = [
                    {
                        "id": m["id"],
                        "name": m.get("name") or m["id"],
                        "vendor": m.get("vendor") or "",
                        "context": m.get("capabilities", {}).get("limits", {}).get("max_context_window_tokens", 128000),
                        "picker": m.get("model_picker_enabled", False),
                    }
                    for m in data.get("data", [])
                ]
                return {"ok": True, "models": models}
            except Exception as exc:
                logger.error("api_copilot_models error: %s", exc)
                return {"ok": False, "error": str(exc), "models": []}

        @app.post("/api/config")
        async def api_save_config(payload: ConfigPayload):
            data = payload.model_dump()
            # upsert named config
            try:
                cfg_blob = _load_configs()
                configs = cfg_blob.get("configs", [])
                name = payload.name or payload.model or f"cfg-{len(configs)+1}"
                # create config dict
                cfg = {
                    "name": name,
                    "provider": payload.provider,
                    "model": payload.model,
                    "openai_api_key": payload.openai_api_key,
                    "openai_base_url": payload.openai_base_url,
                    "anthropic_api_key": payload.anthropic_api_key,
                    "openrouter_api_key": payload.openrouter_api_key,
                    "openrouter_base_url": payload.openrouter_base_url,
                    "github_copilot_token": payload.github_copilot_token,
                    "github_oauth_client_id": payload.github_oauth_client_id,
                    "vllm_base_url": payload.vllm_base_url,
                    "context_window": payload.context_window,
                }
                # replace or append
                found = False
                for i, c in enumerate(configs):
                    if c.get("name") == name:
                        configs[i] = cfg
                        found = True
                        break
                if not found:
                    configs.append(cfg)
                cfg_blob["configs"] = configs
                # optionally set active if payload.name equals active or if no active exists
                if cfg_blob.get("active") is None:
                    cfg_blob["active"] = name
                _save_configs(cfg_blob)

                # If the saved config is the active one, apply immediately
                if cfg_blob.get("active") == name:
                    agent.provider = _provider_from_payload(payload)
                    agent.model = payload.model or agent.model
                    agent.context_window = payload.context_window
                    _runtime_config.update({k: v for k, v in data.items() if v})
                    _persist_env(data)
                return {"ok": True, "name": name}
            except Exception as exc:
                logger.error("api_save_config error: %s", exc)
                return {"ok": False, "error": str(exc)}

        @app.post("/api/config/activate")
        async def api_activate_config(body: dict):
            name = body.get("name")
            if not name:
                return {"ok": False, "error": "name required"}
            cfg_blob = _load_configs()
            configs = cfg_blob.get("configs", [])
            selected = next((c for c in configs if c.get("name") == name), None)
            if not selected:
                return {"ok": False, "error": "config not found"}
            # apply
            try:
                from pydantic import BaseModel
                # use ConfigPayload to build provider
                cp = ConfigPayload(**{
                    "provider": selected.get("provider"),
                    "model": selected.get("model"),
                    "openai_api_key": selected.get("openai_api_key", ""),
                    "openai_base_url": selected.get("openai_base_url", ""),
                    "anthropic_api_key": selected.get("anthropic_api_key", ""),
                    "openrouter_api_key": selected.get("openrouter_api_key", ""),
                    "openrouter_base_url": selected.get("openrouter_base_url", ""),
                    "github_copilot_token": selected.get("github_copilot_token", ""),
                    "github_oauth_client_id": selected.get("github_oauth_client_id", ""),
                    "vllm_base_url": selected.get("vllm_base_url", ""),
                    "context_window": selected.get("context_window", settings.context_window),
                })
                agent.provider = _provider_from_payload(cp)
                agent.model = cp.model or agent.model
                agent.context_window = cp.context_window
                _runtime_config.update({"provider": cp.provider, "model": cp.model, "context_window": cp.context_window})
                cfg_blob["active"] = name
                _save_configs(cfg_blob)
                _persist_env(cp.model_dump())  # persist all fields, not just model+openai_api_key
                return {"ok": True}
            except Exception as exc:
                logger.error("activate error: %s", exc)
                return {"ok": False, "error": str(exc)}

        @app.delete("/api/config/{name}")
        async def api_delete_config(name: str):
            try:
                blob = _load_configs()
                configs = [c for c in blob.get("configs", []) if c.get("name") != name]
                blob["configs"] = configs
                if blob.get("active") == name:
                    blob["active"] = configs[0].get("name") if configs else None
                _save_configs(blob)
                return {"ok": True}
            except Exception as exc:
                return {"ok": False, "error": str(exc)}

        # ==============================================================
        # API — CHAT (REST fallback)
        # ==============================================================

        @app.post("/api/chat")
        async def api_chat(request: Request):
            body       = await request.json()
            message    = body.get("message", "")
            session_id = body.get("session_id")
            if session_id:
                sess = agent.sessions.get_or_create(session_id, channel="web")
            else:
                sess = agent.new_session(channel="web")
            reply = await agent.chat(message, sess)
            return {"reply": reply, "session_id": sess.id}

        # ==============================================================
        # API — SESSIONS
        # ==============================================================

        @app.get("/api/sessions")
        async def api_sessions():
            return [{"id": s.id, "channel": s.channel, "messages": len(s.history)}
                    for s in agent.sessions.list()]

        @app.delete("/api/sessions/{session_id}")
        async def api_del_session(session_id: str):
            agent.sessions.delete(session_id)
            return {"ok": True}

        # ==============================================================
        # API — MEMORY
        # ==============================================================

        @app.get("/api/memory")
        async def api_memory():
            return {
                "long_term": await agent.memory.read_long_term(),
                "today":     await agent.memory.read_daily(),
            }

        @app.post("/api/memory/longterm")
        async def api_save_lt(payload: MemoryPayload):
            await agent.memory.write_long_term(payload.content)
            return {"ok": True}

        @app.get("/api/memory/grep")
        async def api_grep(q: str = ""):
            if not q:
                return {}
            return await agent.memory.grep(q)

        # ==============================================================
        # API — JOBS
        # ==============================================================

        @app.get("/api/jobs")
        async def api_jobs():
            from swarm.scheduler import scheduler
            return scheduler.list_jobs()

        @app.post("/api/jobs")
        async def api_add_job(payload: JobPayload):
            from swarm.scheduler import scheduler

            async def _task(task_text: str = payload.task):
                sess = agent.new_session(channel="scheduler")
                await agent.chat(task_text, sess)

            try:
                scheduler.add_cron(_task, payload.cron, job_id=payload.id,
                                   kwargs={"task_text": payload.task})
                return {"ok": True}
            except Exception as exc:
                return {"ok": False, "error": str(exc)}

        @app.delete("/api/jobs/{job_id}")
        async def api_del_job(job_id: str):
            from swarm.scheduler import scheduler
            try:
                scheduler.remove_job(job_id)
                return {"ok": True}
            except Exception as exc:
                return {"ok": False, "error": str(exc)}

        # ==============================================================
        # API — AGENTS (create / list / docs / chat)
        # ==============================================================

        _agents_dir = Path(__file__).resolve().parent.parent.parent / "agents"
        _sub_agent_cache: dict[str, Any] = {}

        # Inject agents_dir + cache into agent-management tools
        for _tool_name in ("list_agents", "create_agent", "update_agent_instructions"):
            try:
                _t = agent.tools.get(_tool_name)
                _t._agents_dir = _agents_dir
                _t._cache = _sub_agent_cache
            except KeyError:
                pass
        # Disable save_my_instructions for the main Swarm agent (it's only for sub-agents)
        try:
            agent.tools._tools.pop("save_my_instructions", None)
        except Exception:
            pass

        _DEFAULT_AGENT_INSTRUCTIONS = (
            "You are a persistent sub-agent. "
            "You have your own memory and documents. "
            "Be concise and helpful.\n\n"
            "## Auto-configuration\n"
            "You have a special tool called `save_my_instructions`. "
            "Use it whenever the user describes your role, identity, personality, or behaviour — "
            "even partially. Synthesise everything the user has told you about your function "
            "into a single, complete set of instructions and call `save_my_instructions` with "
            "that text. Confirm to the user that your instructions have been saved permanently."
        )

        def _get_sub_agent(name: str):
            """Return (or lazily create) a persistent Agent for this sub-agent."""
            from swarm.core.agent import Agent as _Agent  # local import to avoid circular
            from swarm.memory.manager import MemoryManager as _MM
            if name not in _sub_agent_cache:
                target = _agents_dir / name
                mem = _MM(memory_dir=target / "memory", long_term_file="long_term.txt")
                docs_dir = target / "docs"
                docs_dir.mkdir(exist_ok=True)
                prov = _runtime_config.get("provider") or settings.default_provider
                mdl  = _runtime_config.get("model") or settings.default_model
                # Read saved LLM config
                cfg_path = target / "config.json"
                if cfg_path.exists():
                    saved = json.loads(cfg_path.read_text("utf-8"))
                    prov = saved.get("provider", prov)
                    mdl  = saved.get("model", mdl)
                # Read instructions from agent.md (falls back to default)
                instructions_path = target / "agent.md"
                if instructions_path.exists():
                    instructions = instructions_path.read_text("utf-8").strip() or _DEFAULT_AGENT_INSTRUCTIONS
                else:
                    instructions = _DEFAULT_AGENT_INSTRUCTIONS
                sub = _Agent(
                    provider_name=prov,
                    model=mdl,
                    memory_manager=mem,
                    system_prompt=instructions,
                )
                _sub_agent_cache[name] = sub
            return _sub_agent_cache[name]

        @app.get("/api/agents")
        async def api_list_agents():
            out: list[dict] = []
            try:
                _agents_dir.mkdir(parents=True, exist_ok=True)
                for p in sorted(_agents_dir.iterdir()):
                    if p.is_dir():
                        out.append({"name": p.name})
            except Exception as exc:
                return {"ok": False, "error": str(exc), "agents": []}
            return {"ok": True, "agents": out}

        @app.post("/api/agents")
        async def api_create_agent(body: dict):
            body = body or {}
            name        = body.get("name", "").strip()
            description = body.get("description", "").strip()
            instructions_txt = body.get("instructions", "").strip()
            if not name:
                return {"ok": False, "error": "name required"}
            import re
            if not re.match(r'^[\w\-]+$', name):
                return {"ok": False, "error": "name must only contain letters, digits, hyphens or underscores"}
            try:
                _agents_dir.mkdir(parents=True, exist_ok=True)
                target = _agents_dir / name
                if target.exists():
                    return {"ok": False, "error": "agent already exists"}
                target.mkdir()
                for sub in ("memory", "docs", "sessions"):
                    (target / sub).mkdir()
                (target / "memory" / "long_term.txt").write_text("", encoding="utf-8")
                cfg = {
                    "provider":    _runtime_config.get("provider"),
                    "model":       _runtime_config.get("model"),
                    "description": description,
                }
                (target / "config.json").write_text(json.dumps(cfg, indent=2), encoding="utf-8")
                # Instructions file (agent.md) — use user-provided text or sensible default
                agent_md = instructions_txt if instructions_txt else (
                    f"# Agent: {name}\n\n"
                    "You are a persistent sub-agent with your own memory and documents.\n"
                    "Be concise, focused and helpful.\n"
                )
                (target / "agent.md").write_text(agent_md, encoding="utf-8")
                return {"ok": True, "name": name}
            except Exception as exc:
                return {"ok": False, "error": str(exc)}

        @app.get("/api/agents/{name}/config")
        async def api_agent_config(name: str):
            target = _agents_dir / name
            if not target.is_dir():
                return {"ok": False, "error": "agent not found"}
            cfg_path = target / "config.json"
            cfg = json.loads(cfg_path.read_text("utf-8")) if cfg_path.exists() else {}
            cfg.setdefault("provider", _runtime_config.get("provider"))
            cfg.setdefault("model",    _runtime_config.get("model"))
            return {"ok": True, "config": cfg}

        @app.delete("/api/agents/{name}")
        async def api_delete_agent(name: str):
            import re
            if not re.match(r'^[\w\-]+$', name):
                return {"ok": False, "error": "invalid agent name"}
            target = _agents_dir / name
            if not target.is_dir():
                return {"ok": False, "error": "agent not found"}
            try:
                import shutil
                shutil.rmtree(target)
                _sub_agent_cache.pop(name, None)
                return {"ok": True}
            except Exception as exc:
                return {"ok": False, "error": str(exc)}

        @app.get("/api/agents/{name}/session/{session_id}/history")
        async def api_agent_session_history(name: str, session_id: str):
            """Return the visible (non-system) messages of an agent session."""
            sub = _sub_agent_cache.get(name)
            if not sub:
                return {"ok": True, "messages": []}  # agent not loaded yet — no history
            sess = sub.sessions.get(session_id)
            if not sess:
                return {"ok": True, "messages": []}
            messages = [
                {"role": m.role, "content": m.content}
                for m in sess.history
                if m.role in ("user", "assistant")
            ]
            return {"ok": True, "messages": messages}

        @app.get("/api/agents/{name}/instructions")
        async def api_agent_get_instructions(name: str):
            target = _agents_dir / name
            if not target.is_dir():
                return {"ok": False, "error": "agent not found"}
            instr_path = target / "agent.md"
            content = instr_path.read_text("utf-8") if instr_path.exists() else _DEFAULT_AGENT_INSTRUCTIONS
            return {"ok": True, "content": content}

        @app.put("/api/agents/{name}/instructions")
        async def api_agent_put_instructions(name: str, body: dict):
            target = _agents_dir / name
            if not target.is_dir():
                return {"ok": False, "error": "agent not found"}
            content = (body or {}).get("content", "")
            (target / "agent.md").write_text(content, encoding="utf-8")
            # Invalidate cache so next chat reloads the updated instructions
            _sub_agent_cache.pop(name, None)
            return {"ok": True}

        @app.get("/api/agents/{name}/docs")
        async def api_agent_list_docs(name: str):
            docs_dir = _agents_dir / name / "docs"
            if not docs_dir.is_dir():
                return {"ok": False, "error": "agent not found", "docs": []}
            docs = [{"name": f.name, "size": f.stat().st_size}
                    for f in sorted(docs_dir.iterdir()) if f.is_file()]
            return {"ok": True, "docs": docs}

        @app.post("/api/agents/{name}/docs")
        async def api_agent_add_doc(name: str, body: dict):
            docs_dir = _agents_dir / name / "docs"
            if not docs_dir.is_dir():
                return {"ok": False, "error": "agent not found"}
            filename = (body or {}).get("filename", "").strip()
            content  = (body or {}).get("content", "")
            if not filename:
                return {"ok": False, "error": "filename required"}
            import re
            filename = re.sub(r'[^\w\-\.]+', '_', filename)
            (docs_dir / filename).write_text(content, encoding="utf-8")
            _sub_agent_cache.pop(name, None)
            return {"ok": True, "filename": filename}

        @app.post("/api/agents/{name}/docs/upload")
        async def api_agent_upload_doc(name: str, file: UploadFile = File(...)):
            """Store the file with its original name/format. Text extraction happens at context-injection time."""
            docs_dir = _agents_dir / name / "docs"
            if not docs_dir.is_dir():
                return {"ok": False, "error": "agent not found"}
            import re
            raw  = await file.read()
            orig = re.sub(r'[^\w\-\.]+', '_', file.filename or "document.bin")
            (docs_dir / orig).write_bytes(raw)
            _sub_agent_cache.pop(name, None)
            return {"ok": True, "filename": orig}

        @app.delete("/api/agents/{name}/docs/{doc_name}")
        async def api_agent_del_doc(name: str, doc_name: str):
            doc_path = _agents_dir / name / "docs" / doc_name
            if not doc_path.is_file():
                return {"ok": False, "error": "document not found"}
            doc_path.unlink()
            _sub_agent_cache.pop(name, None)
            return {"ok": True}

        # ==============================================================
        # WebSocket — per sub-agent streaming chat
        # ==============================================================

        def _extract_doc_text(path) -> str:
            """Return the text content of a doc file, extracting from binary formats."""
            from io import BytesIO
            import re
            ext = (re.search(r'\.[^.]+$', path.name.lower()) or type('o',(),{'group':lambda s,n:''})()).group(0)
            raw = path.read_bytes()
            try:
                if ext in ('.pptx', '.ppt'):
                    from pptx import Presentation
                    prs = Presentation(BytesIO(raw))
                    lines = []
                    for i, slide in enumerate(prs.slides, 1):
                        lines.append(f"## Slide {i}")
                        for shape in slide.shapes:
                            if hasattr(shape, 'text') and shape.text.strip():
                                lines.append(shape.text.strip())
                    return '\n'.join(lines)
                if ext in ('.docx', '.doc'):
                    from docx import Document
                    doc = Document(BytesIO(raw))
                    return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
                if ext in ('.xlsx', '.xls'):
                    import openpyxl
                    wb = openpyxl.load_workbook(BytesIO(raw), data_only=True)
                    lines = []
                    for sheet in wb.worksheets:
                        lines.append(f"## Sheet: {sheet.title}")
                        for row in sheet.iter_rows(values_only=True):
                            vals = [str(c) if c is not None else '' for c in row]
                            if any(v.strip() for v in vals):
                                lines.append('\t'.join(vals))
                    return '\n'.join(lines)
                if ext == '.pdf':
                    from pypdf import PdfReader
                    reader = PdfReader(BytesIO(raw))
                    return '\n\n'.join(page.extract_text() or '' for page in reader.pages)
            except Exception as exc:
                return f"[Extraction error for {path.name}: {exc}]"
            # Plain text fallback
            return raw.decode('utf-8', errors='replace')

        @app.websocket("/ws/agents/{name}/{session_id}")
        async def ws_agent_chat(websocket: WebSocket, name: str, session_id: str):
            await websocket.accept()
            target = _agents_dir / name
            if not target.is_dir():
                await websocket.send_json({"type": "error", "message": f"Agent '{name}' not found"})
                await websocket.close()
                return
            sub = _get_sub_agent(name)
            # Inject save_my_instructions tool with this agent's specific path
            try:
                _smi = sub.tools.get("save_my_instructions")
                _smi._agent_md_path = target / "agent.md"
                _smi._agent_name = name
                _smi._cache = _sub_agent_cache
            except KeyError:
                pass
            try:
                existing = sub.sessions.get(session_id)
                sess = existing if existing else sub.new_session(channel="web", session_id=session_id)
            except Exception:
                sess = sub.new_session(channel="web", session_id=session_id)

            agent_docs_dir = target / "docs"

            def _load_docs() -> tuple[str, str]:
                """Return (doc_names_csv, full_docs_ctx) for current docs directory."""
                try:
                    doc_files = [d for d in sorted(agent_docs_dir.iterdir()) if d.is_file()]
                    if not doc_files:
                        return "", ""
                    doc_names = ", ".join(d.name for d in doc_files)
                    docs_ctx = ""
                    for doc in doc_files:
                        text = _extract_doc_text(doc)[:30000]
                        docs_ctx += f"\n\n### Document: {doc.name}\n{text}"
                        logger.info("[docs] Loaded %s — %d chars", doc.name, len(text))
                    return doc_names, docs_ctx
                except Exception as exc:
                    logger.warning("[docs] Load error: %s", exc)
                    return "", ""

            _SELF_CONFIG_META = (
                "\n\n## Auto-configuration\n"
                "Tu disposes d'un outil `save_my_instructions`. "
                "Utilise-le dès que l'utilisateur décrit ou modifie ton rôle, ton identité, ta personnalité "
                "ou ton comportement. Synthétise l'ensemble des instructions dans un texte complet et appelle "
                "`save_my_instructions` avec ce texte. Confirme ensuite à l'utilisateur que ses instructions "
                "ont bien été sauvegardées de façon permanente."
            )

            def _rebuild_system_prompt(doc_names: str, docs_ctx: str) -> None:
                """Rebuild system prompt with docs injected before agent instructions."""
                base_prompt = sub._build_system_context()
                if docs_ctx:
                    header = (
                        f"CONTEXTE — Les documents suivants sont déjà chargés dans ton contexte et leur contenu "
                        f"est disponible ci-dessous. Tu n'as PAS besoin de demander à l'utilisateur de les "
                        f"téléverser — ils sont déjà là. Documents disponibles : {doc_names}\n"
                    )
                    full_prompt = header + "\n## Contenu des documents\n" + docs_ctx + "\n\n## Instructions de l'agent\n" + base_prompt
                else:
                    full_prompt = base_prompt
                # Always append the self-configuration meta-instruction
                full_prompt += _SELF_CONFIG_META
                sess.system_prompt = full_prompt
                logger.info("[docs] System prompt rebuilt — %d chars, docs=%s", len(full_prompt), doc_names or 'none')

            # Initial build on connect
            doc_names, docs_ctx = _load_docs()
            _rebuild_system_prompt(doc_names, docs_ctx)

            try:
                while True:
                    data    = await websocket.receive_json()
                    message = data.get("message", "")
                    if not message:
                        continue
                    # Reload docs before each message so newly uploaded files are visible
                    doc_names, docs_ctx = _load_docs()
                    _rebuild_system_prompt(doc_names, docs_ctx)
                    await websocket.send_json({"type": "thinking_start"})
                    try:
                        async for event in sub.chat_and_stream(message, sess):
                            await websocket.send_json(event)
                        await websocket.send_json({"type": "end"})
                        # context usage stats
                        ctx = sess.token_stats(model=sub.model, limit=sub.context_window)
                        await websocket.send_json({"type": "ctx", **ctx})
                    except Exception as exc:
                        import traceback
                        tb = traceback.format_exc()
                        logger.error("Agent stream error: %s\n%s", exc, tb)
                        await websocket.send_json({"type": "error", "message": str(exc), "detail": tb})
            except WebSocketDisconnect:
                pass

        # ==============================================================
        # WebSocket — streaming chat (global agent)
        # ==============================================================

        @app.websocket("/ws/{session_id}")
        async def ws_chat(websocket: WebSocket, session_id: str):
            await websocket.accept()
            # Use new_session so long-term memory is always injected into the system prompt
            existing = agent.sessions.get(session_id)
            if existing:
                sess = existing
            else:
                sess = agent.new_session(channel="web", session_id=session_id)
            try:
                while True:
                    data    = await websocket.receive_json()
                    message = data.get("message", "")
                    if not message:
                        continue
                    await websocket.send_json({"type": "start"})
                    try:
                        async for chunk in agent.stream_chat(message, sess):
                            await websocket.send_json({"type": "chunk", "text": chunk})
                        await websocket.send_json({"type": "end"})
                        # context usage stats
                        ctx = sess.token_stats(model=agent.model, limit=agent.context_window)
                        await websocket.send_json({"type": "ctx", **ctx})
                    except Exception as exc:
                        import traceback
                        tb = traceback.format_exc()
                        logger.error("Stream error: %s\n%s", exc, tb)
                        await websocket.send_json({
                            "type": "error",
                            "message": str(exc),
                            "detail": tb,
                        })
            except WebSocketDisconnect:
                pass

        return app

    async def start(self) -> None:
        config = uvicorn.Config(
            self.app,
            host=self._host,
            port=self._port,
            log_level="warning",
        )
        self._server = uvicorn.Server(config)
        logger.info("Web gateway -> http://%s:%s", self._host, self._port)
        await self._server.serve()

    async def stop(self) -> None:
        if self._server:
            self._server.should_exit = True
